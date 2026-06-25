# -*- coding: utf-8 -*-
"""
assess_batch_complexity.py — 评估文件复杂度，计算最优批次大小。

复杂度评分公式（追求最高质量，宁可少处理也不超载）:
  基础：每个文件 +1.0
  页数（有文本层）：+0.3/页
  页数（无文本层，需OCR）：+1.0/页
  嵌入图片（需视觉理解）：+2.0/张
  表格：+0.2/个
  大表格（>20格）：额外 +0.5/个
  横向版式（landscape）：+1.0（宽表格/图表，需要更精细的版式还原）

复杂度预算上限：12 分/批（最稳定）、15 分/批（可接受）
- 单文件 ≥ 10 分 → 独占一批（1 个文件）
- 单文件 ≥ 16 分 → 独占一批 + 警告（可能仍需分段处理）

输出 JSON，方便主流程解析。
"""

import sys
import json
import os
import zipfile
from collections import Counter
from datetime import datetime


def image_count_pypdf(page):
    try:
        return len(page.images)
    except Exception:
        pass

    try:
        resources = page.get("/Resources") or {}
        xobjects = resources.get("/XObject") or {}
        count = 0
        for obj in xobjects.values():
            resolved = obj.get_object()
            if resolved.get("/Subtype") == "/Image":
                count += 1
        return count
    except Exception:
        return 0


def inspect_file(filepath):
    """检查单个文件，返回复杂度信息"""
    ext = os.path.splitext(filepath)[1].lower()
    fname = os.path.basename(filepath)
    fsize = os.path.getsize(filepath)
    
    result = {
        "name": fname,
        "path": filepath,
        "type": ext,
        "size_mb": round(fsize / (1024 * 1024), 2),
        "pages": 0,
        "has_text_layer": False,
        "embedded_images": 0,
        "tables": 0,
        "large_tables": 0,
        "orientation": "portrait",  # portrait 或 landscape
        "page_width_pt": 595.0,     # A4 portrait width in points
        "page_height_pt": 842.0,    # A4 portrait height in points
        "score": 0.0,
        "tier": "unknown",
        "notes": []
    }
    
    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(filepath)
            result["pages"] = len(doc)
            
            # 检测页面版式：取所有页面的主流宽度/高度和朝向
            page_orientations = []
            page_widths = []
            page_heights = []
            for page in doc:
                rect = page.rect
                page_widths.append(rect.width)
                page_heights.append(rect.height)
                page_orientations.append("landscape" if rect.width > rect.height else "portrait")
            # 多数页面的朝向决定整体朝向
            result["orientation"] = max(set(page_orientations), key=page_orientations.count)
            result["page_width_pt"] = round(sum(page_widths) / len(page_widths), 1)
            result["page_height_pt"] = round(sum(page_heights) / len(page_heights), 1)
            if result["orientation"] == "landscape":
                result["notes"].append(f"横向版式（{result['page_width_pt']:.0f}×{result['page_height_pt']:.0f}pt）")
            
            # 检查前3页是否有文本层
            text_pages = 0
            for i in range(min(3, len(doc))):
                page_text = doc[i].get_text().strip()
                if len(page_text) > 50:  # 有实质性文本内容
                    text_pages += 1
            
            result["has_text_layer"] = text_pages >= 2  # 至少2/3页有文本
            result["notes"].append(
                f"PDF: {result['pages']}页, "
                f"{'有文本层' if result['has_text_layer'] else '无文本层（需OCR）'}"
            )
            
            # 统计嵌入图片
            img_count = 0
            for page in doc:
                img_count += len(page.get_images())
            result["embedded_images"] = img_count
            if img_count > 0:
                result["notes"].append(f"嵌入图片: {img_count}张")
            
            doc.close()
            
        except ImportError:
            try:
                from pypdf import PdfReader

                reader = PdfReader(filepath)
                result["pages"] = len(reader.pages)
                page_orientations = []
                page_widths = []
                page_heights = []
                text_pages = 0
                img_count = 0

                for page in reader.pages:
                    try:
                        text = (page.extract_text() or "").strip()
                    except Exception:
                        text = ""
                    if len(text) > 50:
                        text_pages += 1

                    img_count += image_count_pypdf(page)
                    try:
                        w, h = float(page.mediabox.width), float(page.mediabox.height)
                    except Exception:
                        w, h = 595.0, 842.0
                    page_widths.append(w)
                    page_heights.append(h)
                    page_orientations.append("landscape" if w > h else "portrait")

                if page_orientations:
                    result["orientation"] = Counter(page_orientations).most_common(1)[0][0]
                    result["page_width_pt"] = round(sum(page_widths) / len(page_widths), 1)
                    result["page_height_pt"] = round(sum(page_heights) / len(page_heights), 1)

                inspected_pages = min(3, max(1, result["pages"]))
                result["has_text_layer"] = text_pages >= min(2, inspected_pages)
                result["embedded_images"] = img_count
                result["notes"].append(
                    f"PDF via pypdf: {result['pages']}页, "
                    f"{'有文本层' if result['has_text_layer'] else '无文本层或文本较少（可能需OCR）'}"
                )
                if img_count > 0:
                    result["notes"].append(f"嵌入图片: {img_count}张")
                if result["orientation"] == "landscape":
                    result["notes"].append(
                        f"横向版式（{result['page_width_pt']:.0f}×{result['page_height_pt']:.0f}pt）"
                    )
            except Exception as e:
                result["notes"].append(f"PDF读取失败: {e}")
                result["pages"] = 1  # 保守估计
        except Exception as e:
            result["notes"].append(f"PDF读取失败: {e}")
            result["pages"] = 1  # 保守估计
    
    elif ext == ".docx":
        try:
            from docx import Document
            from docx.enum.section import WD_ORIENT
            doc = Document(filepath)
            
            # 检测页面版式（取第一个section的朝向）
            if doc.sections:
                section = doc.sections[0]
                # python-docx: section.orientation 返回 WD_ORIENT.LANDSCAPE 或 WD_ORIENT.PORTRAIT
                is_landscape = (section.orientation == WD_ORIENT.LANDSCAPE)
                if is_landscape:
                    result["orientation"] = "landscape"
                    result["page_width_pt"] = 842.0   # A4 landscape
                    result["page_height_pt"] = 595.0
                    result["notes"].append("横向版式")
                # page_width/page_height from section
                result["page_width_pt"] = round(section.page_width.pt, 1) if section.page_width else result["page_width_pt"]
                result["page_height_pt"] = round(section.page_height.pt, 1) if section.page_height else result["page_height_pt"]
            
            # 统计段落（粗略估算页数：每30段 ~ 1页）
            non_empty_paras = [p for p in doc.paragraphs if p.text.strip()]
            para_count = len(non_empty_paras)
            text_chars = sum(len(p.text.strip()) for p in non_empty_paras)
            result["pages"] = max(1, (para_count + 24) // 25)  # 粗略估算
            
            # 统计表格
            result["tables"] = len(doc.tables)
            for table in doc.tables:
                cells = len(table.rows) * len(table.columns)
                if cells > 20:
                    result["large_tables"] += 1
            
            # 检查嵌入图片
            with zipfile.ZipFile(filepath, 'r') as z:
                media_files = [f for f in z.namelist() if 'word/media/' in f]
                result["embedded_images"] = len(media_files)
            
            if result["embedded_images"] > 0:
                result["notes"].append(f"嵌入图片: {result['embedded_images']}张")
            if result["tables"] > 0:
                result["notes"].append(f"表格: {result['tables']}个（大表: {result['large_tables']}）")
            if result["embedded_images"] > len(doc.paragraphs) * 0.5:
                # 图片型DOCX：页面主要由图片组成
                result["pages"] = result["embedded_images"]
                result["has_text_layer"] = False
                result["notes"].append("图片型DOCX，每张图片视作一页")
            else:
                result["has_text_layer"] = text_chars > 0
                result["notes"].append(
                    f"DOCX: ~{result['pages']}页, "
                    f"{'有可读文本' if result['has_text_layer'] else '文本较少'}"
                )
                
        except Exception as e:
            result["notes"].append(f"DOCX读取失败: {e}")
            result["pages"] = max(1, result["embedded_images"])
    
    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            result["pages"] = len(prs.slides)
            result["has_text_layer"] = True
            result["notes"].append(f"PPT: {result['pages']}张幻灯片")
            # 统计嵌入图片
            img_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        img_count += 1
            result["embedded_images"] = img_count
        except Exception as e:
            result["notes"].append(f"PPT读取失败: {e}")
            result["pages"] = 10  # 保守估计
    
    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            result["pages"] = len(wb.sheetnames)
            # 估算表格复杂度
            total_cells = 0
            for name in wb.sheetnames:
                ws = wb[name]
                total_cells += ws.max_row * ws.max_column
            result["notes"].append(f"Excel: {len(wb.sheetnames)}张工作表, ~{total_cells}格")
            # 大表格：折算为复杂度
            if total_cells > 100:
                result["tables"] = len(wb.sheetnames)
                result["large_tables"] = len(wb.sheetnames)
            wb.close()
        except Exception as e:
            result["notes"].append(f"Excel读取失败: {e}")
            result["pages"] = 1
    
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"):
        result["pages"] = 1
        result["has_text_layer"] = False
        result["embedded_images"] = 1
        result["notes"].append(f"图片文件: 1张（需OCR）")
    
    else:
        result["pages"] = 1
        result["notes"].append(f"未知格式: {ext}")
    
    # === 计算复杂度分数 ===
    score = 1.0  # 基础分
    
    if result["has_text_layer"]:
        score += result["pages"] * 0.3
    else:
        score += result["pages"] * 1.0  # OCR页更贵
    
    score += result["embedded_images"] * 2.0
    score += result["tables"] * 0.2
    score += result["large_tables"] * 0.5
    
    # 横向版式需要更精细的版式还原（宽表格、多列布局），增加复杂度
    if result["orientation"] == "landscape":
        score += 1.0
    
    result["score"] = round(score, 1)
    
    # === 复杂度分级 ===
    if score < 4:
        result["tier"] = "simple"
    elif score < 8:
        result["tier"] = "medium"
    elif score < 14:
        result["tier"] = "complex"
    else:
        result["tier"] = "very-complex"
    
    return result


def build_batches(files_info, budget=12.0):
    """
    贪心算法构建批次：逐个添加文件，累计分不超过预算。
    单个文件超过预算的独占一批。
    """
    batches = []
    current_batch = []
    current_score = 0.0
    
    for fi in files_info:
        if fi["score"] >= budget:
            # 这个文件本身就超预算——独占一批
            if current_batch:
                batches.append({
                    "files": current_batch,
                    "total_score": round(current_score, 1),
                    "count": len(current_batch)
                })
                current_batch = []
                current_score = 0.0
            batches.append({
                "files": [fi],
                "total_score": fi["score"],
                "count": 1
            })
        elif current_score + fi["score"] <= budget:
            current_batch.append(fi)
            current_score += fi["score"]
        else:
            # 当前批次满了，结算
            batches.append({
                "files": current_batch,
                "total_score": round(current_score, 1),
                "count": len(current_batch)
            })
            current_batch = [fi]
            current_score = fi["score"]
    
    # 最后一个批次
    if current_batch:
        batches.append({
            "files": current_batch,
            "total_score": round(current_score, 1),
            "count": len(current_batch)
        })
    
    return batches


def main():
    import argparse
    parser = argparse.ArgumentParser(description="评估文件复杂度并分批")
    parser.add_argument("files", nargs="*", help="待处理文件列表")
    parser.add_argument("--tbd-dir", help="TBD目录路径，自动扫描")
    parser.add_argument("--budget", type=float, default=12.0, help="每批复杂度预算上限（默认12）")
    parser.add_argument("--json", action="store_true", help="输出JSON格式")
    args = parser.parse_args()
    
    files = []
    
    if args.tbd_dir:
        tbd = args.tbd_dir
        if os.path.isdir(tbd):
            for f in sorted(os.listdir(tbd)):
                fpath = os.path.join(tbd, f)
                if os.path.isfile(fpath) and not f.startswith("."):
                    files.append(fpath)
    
    for f in args.files:
        if f not in files:
            files.append(f)
    
    if not files:
        if args.json:
            print(json.dumps({"error": "no files found", "batches": []}, ensure_ascii=False))
        else:
            print("没有找到文件。")
        return 0
    
    # 检查每个文件
    all_info = []
    for f in files:
        info = inspect_file(f)
        all_info.append(info)
    
    # 构建批次
    batches = build_batches(all_info, args.budget)
    
    if args.json:
        output = {
            "total_files": len(all_info),
            "total_batches": len(batches),
            "budget_per_batch": args.budget,
            "files": all_info,
            "batches": batches
        }
        print(json.dumps(output, ensure_ascii=False, indent=2))
    else:
        print(f"评估 {len(all_info)} 个文件，复杂度预算 {args.budget} 分/批")
        print(f"分为 {len(batches)} 批处理\n")
        print(f"{'文件名':<45} {'类型':<8} {'页数':>4} {'图片':>4} {'版式':>5} {'分数':>6} {'级别':<12}")
        print("-" * 95)
        
        for fi in all_info:
            icon = {"simple": "🟢", "medium": "🟡", "complex": "🟠", "very-complex": "🔴"}.get(fi["tier"], "⚪")
            name = fi["name"][:42] + "..." if len(fi["name"]) > 45 else fi["name"]
            orient_label = "横" if fi["orientation"] == "landscape" else "纵"
            print(f"{name:<45} {fi['type']:<8} {fi['pages']:>4} {fi['embedded_images']:>4} {orient_label:>5} {fi['score']:>6.1f} {icon} {fi['tier']:<10}")
        
        print()
        for i, batch in enumerate(batches, 1):
            names = [f["name"] for f in batch["files"]]
            print(f"第 {i} 批（{batch['count']} 个文件，总分 {batch['total_score']}）：{' | '.join(names)}")
        
        # 警告
        for fi in all_info:
            if fi["tier"] == "very-complex":
                print(f"\n⚠️ 警告：'{fi['name']}' 复杂度极高（{fi['score']}分），独占一批，仍可能影响质量。")
    
    return 0


if __name__ == "__main__":
    sys.exit(main())
